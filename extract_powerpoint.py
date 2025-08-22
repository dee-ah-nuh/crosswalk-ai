#!/usr/bin/env python3
"""
Extract content from PowerPoint presentation about the data model
"""

try:
    from pptx import Presentation
    from pptx.util import Inches
    import os
    
    def extract_powerpoint_content(file_path):
        """Extract text content from PowerPoint presentation"""
        print(f"Extracting content from: {file_path}")
        
        # Load the presentation
        prs = Presentation(file_path)
        
        print(f"Total slides: {len(prs.slides)}")
        
        all_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"\n=== SLIDE {slide_num} ===")
            slide_content = []
            
            # Extract text from all text boxes and shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    print(f"Text: {text}")
                    slide_content.append(text)
                
                # Check for tables
                if hasattr(shape, 'table'):
                    print("Found table:")
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        row_data = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_data.append(cell.text.strip())
                        if row_data:
                            print(f"  Row {row_idx}: {' | '.join(row_data)}")
                            slide_content.append(' | '.join(row_data))
            
            all_content.append({
                'slide': slide_num,
                'content': slide_content
            })
        
        return all_content
    
    # Extract the PowerPoint content
    ppt_file = "attached_assets/New_Data_Model_Training_20240313v1_1755827391509.pptx"
    
    if os.path.exists(ppt_file):
        content = extract_powerpoint_content(ppt_file)
        
        # Save extracted content to text file for analysis
        with open("data_model_content.txt", "w", encoding="utf-8") as f:
            f.write("EXTRACTED DATA MODEL TRAINING CONTENT\n")
            f.write("=" * 50 + "\n\n")
            
            for slide_info in content:
                f.write(f"SLIDE {slide_info['slide']}:\n")
                for text in slide_info['content']:
                    f.write(f"  {text}\n")
                f.write("\n")
        
        print(f"\n✅ Content extracted and saved to data_model_content.txt")
        print(f"Found {len(content)} slides with content")
        
    else:
        print(f"❌ File not found: {ppt_file}")
        
except ImportError:
    print("Installing python-pptx library...")
    import subprocess
    import sys
    
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        print("✅ Installed python-pptx, please run the script again")
    except Exception as e:
        print(f"❌ Failed to install python-pptx: {e}")
        
except Exception as e:
    print(f"❌ Error extracting PowerPoint content: {e}")